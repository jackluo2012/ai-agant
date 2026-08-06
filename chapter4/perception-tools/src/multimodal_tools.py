"""
多模态理解工具模块

提供网页、文档、图片和视频的内容解析功能。

此模块包含：
- read_webpage: 读取并提取网页内容
- read_document: 读取文档内容（PDF、DOCX、PPTX）
- parse_image: 解析图片基本信息
- parse_video: 解析视频元数据
- extract_youtube_transcript: 提取 YouTube 字幕
- download_youtube_video: 下载 YouTube 视频
"""
import json
import logging
import os
import sys
import traceback
from pathlib import Path
from typing import Union
import base64

import requests
from bs4 import BeautifulSoup
from dotenv import load_dotenv
from mcp.types import TextContent
from pydantic import Field

# 添加项目根目录到路径，用于导入统一 LLM 客户端
_project_root = os.path.abspath(os.path.join(os.path.dirname(__file__), "../.."))
if _project_root not in sys.path:
    sys.path.insert(0, _project_root)

from base import ActionResponse, validate_file_path, download_file_from_url, is_url


load_dotenv()


async def read_webpage(
    url: str,
    extract_text: bool = True,
    extract_links: bool = False
) -> Union[str, TextContent]:
    """
    读取并从网页提取内容

    使用 BeautifulSoup 解析网页，提取文本内容和链接。

    Args:
        url: 网页 URL
        extract_text: 是否提取主要文本内容（默认 True）
        extract_links: 是否提取所有链接（默认 False）

    Returns:
        包含提取的网页内容的 TextContent

    示例:
        >>> result = await read_webpage(
        ...     "https://example.com",
        ...     extract_text=True,
        ...     extract_links=True
        ... )
        >>> result['title']
        '网页标题'
    """
    try:
        logging.info(f"📄 正在读取网页：{url}")

        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"
        }

        response = requests.get(url, headers=headers, timeout=30)
        response.raise_for_status()

        soup = BeautifulSoup(response.content, 'html.parser')

        result = {
            "url": url,
            "title": soup.title.string if soup.title else "无标题"
        }

        # 提取文本内容
        if extract_text:
            # 移除脚本和样式元素
            for script in soup(["script", "style", "noscript"]):
                script.decompose()

            text = soup.get_text()
            # 清理空白
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = ' '.join(chunk for chunk in chunks if chunk)

            result["text"] = text[:10000]  # 限制为前 10000 个字符
            result["text_length"] = len(text)

        # 提取链接
        if extract_links:
            links = []
            for link in soup.find_all('a', href=True):
                href = link['href']
                # 补全相对链接
                if href.startswith('/'):
                    from urllib.parse import urljoin
                    href = urljoin(url, href)
                links.append({
                    "text": link.get_text().strip(),
                    "href": href
                })
            result["links"] = links[:100]  # 限制为前 100 个链接

        logging.info(f"✅ 成功提取网页内容")

        action_response = ActionResponse(
            success=True,
            message=result,
            metadata={"url": url}
        )

        return TextContent(
            type="text",
            text=json.dumps(action_response.model_dump(), ensure_ascii=False)
        )

    except Exception as e:
        error_msg = f"网页读取失败：{str(e)}"
        logging.error(f"网页错误：{traceback.format_exc()}")

        action_response = ActionResponse(
            success=False,
            message=error_msg,
            metadata={"error_type": "webpage_error", "url": url}
        )

        return TextContent(
            type="text",
            text=json.dumps(action_response.model_dump(), ensure_ascii=False)
        )


def _sniff_document_type(path: Path) -> str | None:
    """
    通过魔术字节检测文档类型（PDF/DOCX/PPTX）

    当文件扩展名不可靠时（如从 URL 下载的临时文件），
    通过读取文件头部的魔术字节来判断真实文件类型。

    Args:
        path: 文件路径

    Returns:
        文件扩展名（如 '.pdf'、'.docx'）或 None

    注意:
        - PDF: 以 '%PDF' 开头
        - DOCX/PPTX: ZIP 格式（PK\\x03\\x04），需进一步检查内部结构
    """
    try:
        with open(path, 'rb') as f:
            header = f.read(4)

        # PDF 检测
        if header.startswith(b'%PDF'):
            return '.pdf'

        # ZIP 格式（DOCX/PPTX 都是 ZIP）
        if header.startswith(b'PK\x03\x04'):
            import zipfile
            with zipfile.ZipFile(path) as zf:
                names = zf.namelist()

            # 检查是否是 DOCX
            if any(n.startswith('word/') for n in names):
                return '.docx'

            # 检查是否是 PPTX
            if any(n.startswith('ppt/') for n in names):
                return '.pptx'

    except Exception as e:
        logging.debug(f"文档类型检测失败：{e}")

    return None


async def read_document(
    file_path: str,
    extract_images: bool = False
) -> Union[str, TextContent]:
    """
    读取并从文档中提取内容（PDF、DOCX、PPTX）

    支持从本地文件或 URL 读取文档内容。

    Args:
        file_path: 文档文件路径或 URL
        extract_images: 是否从文档中提取图片（暂未实现）

    Returns:
        包含提取的文档内容的 TextContent

    示例:
        >>> result = await read_document("/path/to/document.pdf")
        >>> result['file_type']
        'pdf'
        >>> result['text'][:100]
        '文档内容...'
    """
    try:
        # 处理 URL 下载
        if is_url(file_path):
            logging.info(f"📥 正在从 URL 下载文档")
            temp_path, _ = download_file_from_url(file_path)
            file_path = temp_path

        path = validate_file_path(file_path)

        logging.info(f"📄 正在读取文档：{path}")

        file_ext = path.suffix.lower()
        if file_ext not in ('.pdf', '.docx', '.pptx'):
            # URL 的路径通常没有真正的扩展名
            # 因此通过文件的魔术字节检测真实类型
            file_ext = _sniff_document_type(path) or file_ext

        result = {}

        # PDF 提取
        if file_ext == '.pdf':
            import PyPDF2

            with open(path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                text = ""
                for page in reader.pages:
                    text += page.extract_text() + "\n"

                result = {
                    "file_name": path.name,
                    "file_type": "pdf",
                    "page_count": len(reader.pages),
                    "text": text[:15000],  # 限制大小
                    "text_length": len(text)
                }

        # DOCX 提取
        elif file_ext == '.docx':
            from docx import Document

            doc = Document(path)
            text = "\n".join([para.text for para in doc.paragraphs])

            result = {
                "file_name": path.name,
                "file_type": "docx",
                "paragraph_count": len(doc.paragraphs),
                "text": text[:15000],
                "text_length": len(text)
            }

        # PPTX 提取
        elif file_ext == '.pptx':
            from pptx import Presentation

            prs = Presentation(path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"

            result = {
                "file_name": path.name,
                "file_type": "pptx",
                "slide_count": len(prs.slides),
                "text": text[:15000],
                "text_length": len(text)
            }

        else:
            raise ValueError(f"不支持的文件类型：{file_ext}")

        logging.info(f"✅ 成功提取文档内容")

        action_response = ActionResponse(
            success=True,
            message=result,
            metadata={"file_path": str(path), "file_type": file_ext}
        )

        return TextContent(
            type="text",
            text=json.dumps(action_response.model_dump(), ensure_ascii=False)
        )

    except Exception as e:
        error_msg = f"文档读取失败：{str(e)}"
        logging.error(f"文档错误：{traceback.format_exc()}")

        action_response = ActionResponse(
            success=False,
            message=error_msg,
            metadata={"error_type": "document_error"}
        )

        return TextContent(
            type="text",
            text=json.dumps(action_response.model_dump(), ensure_ascii=False)
        )


async def parse_image(
    image_path: str,
    use_llm: bool = True
) -> Union[str, TextContent]:
    """
    解析和理解图片内容

    提取图片的基本信息，可选进行 LLM 视觉分析。

    Args:
        image_path: 图片文件路径或 URL
        use_llm: 是否使用 LLM 进行图片理解（默认 True）

    Returns:
        包含图片分析的 TextContent

    注意:
        - use_llm=True 时需要支持视觉的 LLM 模型
        - 返回 base64 编码的图片数据可用于视觉 API

    示例:
        >>> result = await parse_image("/path/to/image.jpg", use_llm=False)
        >>> result['format']
        'JPEG'
        >>> result['size']
        (1920, 1080)
    """
    try:
        # 处理 URL 下载
        if is_url(image_path):
            logging.info(f"📥 正在从 URL 下载图片")
            temp_path, _ = download_file_from_url(image_path)
            image_path = temp_path

        path = validate_file_path(image_path)

        logging.info(f"🖼️ 正在解析图片：{path}")

        from PIL import Image

        img = Image.open(path)

        result = {
            "file_name": path.name,
            "format": img.format,
            "mode": img.mode,
            "size": img.size,
            "width": img.width,
            "height": img.height
        }

        # 如果请求 LLM 分析，为视觉 API 编码图片
        if use_llm:
            with open(path, 'rb') as img_file:
                img_base64 = base64.b64encode(img_file.read()).decode('utf-8')
                result["base64_preview"] = img_base64[:100] + "..."  # 预览
                result["note"] = "完整的 base64 数据可用于视觉 API 分析"

        logging.info(f"✅ 成功解析图片")

        action_response = ActionResponse(
            success=True,
            message=result,
            metadata={"file_path": str(path)}
        )

        return TextContent(
            type="text",
            text=json.dumps(action_response.model_dump(), ensure_ascii=False)
        )

    except Exception as e:
        error_msg = f"图片解析失败：{str(e)}"
        logging.error(f"图片错误：{traceback.format_exc()}")

        action_response = ActionResponse(
            success=False,
            message=error_msg,
            metadata={"error_type": "image_error"}
        )

        return TextContent(
            type="text",
            text=json.dumps(action_response.model_dump(), ensure_ascii=False)
        )


async def parse_video(
    video_path: str,
    extract_frames: bool = False,
    frame_interval: int = 30
) -> Union[str, TextContent]:
    """
    解析并从视频文件提取元数据

    提取视频的基本信息，可选提取关键帧。

    Args:
        video_path: 视频文件路径或 URL
        extract_frames: 是否提取示例帧（默认 False）
        frame_interval: 每隔 N 秒提取一帧（默认 30）

    Returns:
        包含视频元数据的 TextContent

    注意:
        - extract_frames=True 时会返回帧的时间戳信息
        - 支持本地文件和 URL 下载

    示例:
        >>> result = await parse_video("/path/to/video.mp4")
        >>> result['duration_seconds']
        120.5
        >>> result['resolution']
        '1920x1080'
    """
    try:
        # 处理 URL 下载
        if is_url(video_path):
            logging.info(f"📥 正在从 URL 下载视频")
            temp_path, _ = download_file_from_url(video_path, max_size_mb=500)
            video_path = temp_path

        path = validate_file_path(video_path)

        logging.info(f"🎥 正在解析视频：{path}")

        import cv2

        video = cv2.VideoCapture(str(path))

        fps = video.get(cv2.CAP_PROP_FPS)
        frame_count = int(video.get(cv2.CAP_PROP_FRAME_COUNT))
        width = int(video.get(cv2.CAP_PROP_FRAME_WIDTH))
        height = int(video.get(cv2.CAP_PROP_FRAME_HEIGHT))
        duration = frame_count / fps if fps > 0 else 0

        result = {
            "file_name": path.name,
            "duration_seconds": round(duration, 2),
            "fps": fps,
            "frame_count": frame_count,
            "resolution": f"{width}x{height}",
            "width": width,
            "height": height
        }

        # 提取关键帧
        if extract_frames:
            frames = []
            interval_frames = int(fps * frame_interval) if fps > 0 else 30

            for frame_num in range(0, frame_count, interval_frames):
                video.set(cv2.CAP_PROP_POS_FRAMES, frame_num)
                ret, frame = video.read()
                if ret:
                    timestamp = frame_num / fps if fps > 0 else 0
                    frames.append({
                        "frame_number": frame_num,
                        "timestamp": round(timestamp, 2)
                    })

            result["keyframes"] = frames[:20]  # 限制为 20 帧
            result["keyframes_count"] = len(frames)

        video.release()

        logging.info(f"✅ 成功解析视频元数据")

        action_response = ActionResponse(
            success=True,
            message=result,
            metadata={"file_path": str(path)}
        )

        return TextContent(
            type="text",
            text=json.dumps(action_response.model_dump(), ensure_ascii=False)
        )

    except Exception as e:
        error_msg = f"视频解析失败：{str(e)}"
        logging.error(f"视频错误：{traceback.format_exc()}")

        action_response = ActionResponse(
            success=False,
            message=error_msg,
            metadata={"error_type": "video_error"}
        )

        return TextContent(
            type="text",
            text=json.dumps(action_response.model_dump(), ensure_ascii=False)
        )


async def download_youtube_video(
    url: str,
    output_dir: str = ".",
    max_resolution: str = "720p"
) -> Union[str, TextContent]:
    """
    使用 yt-dlp 下载 YouTube 视频

    Args:
        url: YouTube 视频 URL
        output_dir: 保存视频的目录（默认当前目录）
        max_resolution: 最大分辨率（360p、480p、720p、1080p）

    Returns:
        包含下载结果的 TextContent

    注意:
        - 需要安装 yt-dlp：pip install yt-dlp
        - 下载的文件名使用视频标题

    示例:
        >>> result = await download_youtube_video(
        ...     "https://youtube.com/watch?v=xxx",
        ...     output_dir="/tmp",
        ...     max_resolution="720p"
        ... )
    """
    try:
        logging.info(f"📥 正在下载 YouTube 视频：{url}")

        try:
            import yt_dlp

            output_template = Path(output_dir) / '%(title)s.%(ext)s'

            ydl_opts = {
                'format': f'best[height<={max_resolution[:-1]}]',
                'outtmpl': str(output_template),
                'quiet': False
            }

            with yt_dlp.YoutubeDL(ydl_opts) as ydl:
                info = ydl.extract_info(url, download=True)

                result = {
                    "title": info['title'],
                    "duration": info.get('duration'),
                    "output_dir": output_dir,
                    "resolution": max_resolution,
                    "video_id": info['id']
                }

                logging.info(f"✅ 已下载：{info['title']}")

                action_response = ActionResponse(
                    success=True,
                    message=result,
                    metadata={"url": url}
                )

        except ImportError:
            raise ImportError(
                "未安装 yt-dlp。"
                "使用以下命令安装：pip install yt-dlp"
            )

        return TextContent(
            type="text",
            text=json.dumps(action_response.model_dump(), ensure_ascii=False)
        )

    except Exception as e:
        error_msg = f"YouTube 下载失败：{str(e)}"
        logging.error(f"YouTube 下载错误：{traceback.format_exc()}")

        action_response = ActionResponse(
            success=False,
            message=error_msg,
            metadata={"error_type": "youtube_download_error"}
        )

        return TextContent(
            type="text",
            text=json.dumps(action_response.model_dump(), ensure_ascii=False)
        )


async def extract_youtube_transcript(
    video_id: str,
    language_code: str = "zh",
    translate_to_language: str | None = None
) -> Union[str, TextContent]:
    """
    从 YouTube 视频提取字幕

    Args:
        video_id: YouTube 视频 ID 或完整 URL
        language_code: 字幕语言代码（默认：zh 中文）
        translate_to_language: 翻译到此语言（可选）

    Returns:
        包含字幕数据的 TextContent

    注意:
        - 需要安装 youtube-transcript-api
        - 支持自动生成的字幕

    示例:
        >>> result = await extract_youtube_transcript(
        ...     "https://youtube.com/watch?v=xxx",
        ...     language_code="zh"
        ... )
        >>> result['transcript'][:3]
        [{'timestamp': '00:00', 'text': '第一句话'}, ...]
    """
    try:
        from youtube_transcript_api import YouTubeTranscriptApi

        # 如果提供了完整 URL，清理 video_id
        if "youtube.com" in video_id or "youtu.be" in video_id:
            if "?v=" in video_id:
                video_id = video_id.split("?v=")[-1].split("&")[0]
            elif "youtu.be/" in video_id:
                video_id = video_id.split("youtu.be/")[-1].split("?")[0]

        logging.info(f"📺 正在提取视频字幕，视频 ID：{video_id}")

        # 使用正确的 API 获取字幕
        if translate_to_language:
            transcript_list = YouTubeTranscriptApi().list(video_id)
            try:
                transcript = transcript_list.find_transcript([language_code])
            except Exception:
                # 如果未找到指定语言，获取任何可用的字幕
                transcript = transcript_list.find_generated_transcript(["zh"])
            # 翻译为目标语言
            fetched_transcript = transcript.translate(translate_to_language).fetch()
            transcript_data = fetched_transcript.snippets
        else:
            try:
                # 使用 fetch 方法获取字幕
                fetched_transcript = YouTubeTranscriptApi().fetch(
                    video_id,
                    languages=(language_code,)
                )
                transcript_data = fetched_transcript.snippets
            except Exception:
                # 回退到中文
                fetched_transcript = YouTubeTranscriptApi().fetch(
                    video_id,
                    languages=("zh",)
                )
                transcript_data = fetched_transcript.snippets

        # 格式化字幕
        formatted_transcript = []
        for entry in transcript_data:
            # 作为对象属性访问
            start_time = entry.start if hasattr(entry, 'start') else entry.get('start', 0)
            text = entry.text if hasattr(entry, 'text') else entry.get('text', '')
            minutes, seconds = divmod(int(start_time), 60)
            timestamp = f"{minutes:02d}:{seconds:02d}"
            formatted_transcript.append({
                "timestamp": timestamp,
                "text": text
            })

        # 创建完整文本版本
        full_text = " ".join([
            entry.text if hasattr(entry, 'text') else entry.get('text', '')
            for entry in transcript_data
        ])

        result = {
            "video_id": video_id,
            "language": translate_to_language if translate_to_language else language_code,
            "transcript": formatted_transcript[:200],  # 限制为前 200 条
            "total_entries": len(transcript_data),
            "full_text": full_text[:10000],  # 限制全文为 10000 个字符
            "full_text_length": len(full_text)
        }

        logging.info(f"✅ 成功提取字幕（{len(transcript_data)} 条）")

        action_response = ActionResponse(
            success=True,
            message=result,
            metadata={
                "video_id": video_id,
                "language": language_code,
                "translated": translate_to_language is not None
            }
        )

        return TextContent(
            type="text",
            text=json.dumps(action_response.model_dump(), ensure_ascii=False)
        )

    except Exception as e:
        error_msg = f"YouTube 字幕提取失败：{str(e)}"
        logging.error(f"YouTube 错误：{traceback.format_exc()}")

        action_response = ActionResponse(
            success=False,
            message=error_msg,
            metadata={"error_type": "youtube_error", "video_id": video_id}
        )

        return TextContent(
            type="text",
            text=json.dumps(action_response.model_dump(), ensure_ascii=False)
        )
